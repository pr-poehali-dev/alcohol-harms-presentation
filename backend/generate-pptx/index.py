"""Генерация PPTX-файла презентации об алкоголе и здоровье (9 слайдов)."""
import base64
import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


BG    = RGBColor(0x0D, 0x0F, 0x14)
RED   = RGBColor(0xC8, 0x33, 0x2A)
LIGHT = RGBColor(0xE8, 0xE0, 0xD5)
GRAY  = RGBColor(0x9A, 0x8F, 0x8A)
DIM   = RGBColor(0x6A, 0x5F, 0x5F)
GREEN = RGBColor(0x3A, 0xAA, 0x70)
BLUE  = RGBColor(0x4A, 0x8C, 0xBF)
MUTED = RGBColor(0xC8, 0xC0, 0xB8)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_label(slide, text, W):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.25), W - Inches(1.2), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text.upper()
    r.font.size = Pt(8)
    r.font.bold = True
    r.font.color.rgb = RED


def add_title(slide, text, W, top=0.65):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), W - Inches(1.2), Inches(1.1))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = text
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = LIGHT


def add_bullets(slide, bullets, W, top=1.85, font_size=14):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(top), W - Inches(1.2), Inches(7.5 - top - 0.4))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for (emoji, title, text) in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(5)
        if emoji:
            r0 = p.add_run()
            r0.text = emoji + "  "
            r0.font.size = Pt(font_size)
            r0.font.color.rgb = MUTED
        if title:
            r1 = p.add_run()
            r1.text = title + ": "
            r1.font.size = Pt(font_size)
            r1.font.bold = True
            r1.font.color.rgb = LIGHT
        r2 = p.add_run()
        r2.text = text
        r2.font.size = Pt(font_size)
        r2.font.color.rgb = MUTED


def add_note_box(slide, text, W, top):
    box = slide.shapes.add_textbox(Inches(0.6), Inches(top), W - Inches(1.2), Inches(0.55))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "⚠  " + text
    r.font.size = Pt(12)
    r.font.color.rgb = GRAY


# ── Слайды ──────────────────────────────────────────────

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    W, H = prs.slide_width, prs.slide_height

    tb = slide.shapes.add_textbox(Inches(1), Inches(1.6), W - Inches(2), Inches(1.5))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Алкоголь — угроза здоровью"
    r.font.size = Pt(40)
    r.font.bold = True
    r.font.color.rgb = LIGHT

    tb2 = slide.shapes.add_textbox(Inches(1), Inches(3.2), W - Inches(2), Inches(0.5))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "Факты, статистика и методы профилактики"
    r2.font.size = Pt(17)
    r2.font.color.rgb = GRAY

    line = slide.shapes.add_shape(1, Inches(3), Inches(4.0), Inches(4), Pt(2))
    line.fill.solid()
    line.fill.fore_color.rgb = RED
    line.line.fill.background()

    tb3 = slide.shapes.add_textbox(Inches(1), Inches(4.3), W - Inches(2), Inches(0.9))
    tf3 = tb3.text_frame
    p3 = tf3.paragraphs[0]
    p3.alignment = PP_ALIGN.CENTER
    r3 = p3.add_run()
    r3.text = "Выполнил: Илья Садомов, 9Д"
    r3.font.size = Pt(14)
    r3.font.color.rgb = LIGHT
    p4 = tf3.add_paragraph()
    p4.alignment = PP_ALIGN.CENTER
    r4 = p4.add_run()
    r4.text = "Руководитель: Екатерина Геннадьевна"
    r4.font.size = Pt(14)
    r4.font.color.rgb = GRAY


def slide_intro(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    W = prs.slide_width
    add_label(slide, "Слайд 1 — Введение", W)
    add_title(slide, "Что такое алкоголь?", W)
    add_bullets(slide, [
        ("🍷", "Определение", "Этанол (C₂H₅OH) — психоактивное вещество, подавляющее ЦНС и вызывающее физическую и психологическую зависимость"),
        ("⚠️", "Классификация ВОЗ", "Алкоголь признан канцерогеном 1-й группы — наравне с табаком, асбестом и радиацией. Безопасной дозы не существует"),
        ("🌍", "Мировая проблема", "Ежегодно от алкоголя погибает более 3 млн человек — больше, чем от СПИДа и туберкулёза вместе взятых"),
        ("🇷🇺", "Россия", "Входит в топ-10 стран по потреблению. Средний возраст первого употребления — 13–14 лет, один из самых низких в Европе"),
        ("🎯", "Цель", "Донести научно обоснованную информацию о вреде алкоголя и показать эффективные методы профилактики"),
    ], W)


def slide_stats_world(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    W = prs.slide_width
    add_label(slide, "Слайд 2 — Мировая статистика", W)
    add_title(slide, "Масштаб проблемы в мире", W)
    add_bullets(slide, [
        ("📊", "3 млн+ смертей в год", "алкоголь — третья причина смертности в мире"),
        ("📊", "5,1% от всех смертей", "прямо или косвенно связаны с употреблением алкоголя"),
        ("📊", "200+ заболеваний", "напрямую вызываются или усугубляются алкоголем"),
        ("📊", "237 млн человек", "в мире страдают расстройствами, связанными с алкоголем"),
        ("🚗", "43% смертельных ДТП", "происходят с участием пьяных водителей"),
        ("💔", "30% самоубийств", "и 40% случаев бытового насилия связаны с алкоголем"),
        ("💰", "1,7 трлн $ в год", "составляет экономический ущерб от алкоголя в мире"),
    ], W)


def slide_stats_russia(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    W = prs.slide_width
    add_label(slide, "Слайд 3 — Статистика в России", W)
    add_title(slide, "Алкоголь в России: цифры Росстата", W)
    add_bullets(slide, [
        ("📈", "70%", "взрослых россиян употребляют алкоголь"),
        ("📈", "35%", "употребляют регулярно — раз в неделю и чаще"),
        ("📈", "8%", "страдают алкогольной зависимостью"),
        ("📈", "52%", "начали пить до 15 лет"),
        ("📉", "На 15–20 лет", "сокращается средняя продолжительность жизни зависимых мужчин"),
        ("🏥", "Каждая 5-я госпитализация", "в России связана с алкоголем"),
        ("👶", "50 000 детей ежегодно", "рождаются с признаками фетального алкогольного синдрома"),
        ("👮", "60%+ преступлений", "в России совершается в состоянии алкогольного опьянения"),
    ], W, font_size=13)


def slide_health_brain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    W = prs.slide_width
    add_label(slide, "Слайд 4 — Влияние на нервную систему", W)
    add_title(slide, "Что алкоголь делает с мозгом", W)
    add_bullets(slide, [
        ("🧠", "Нейроны", "Этанол разрушает нейронные связи. Даже одна «пьяная» ночь уничтожает тысячи клеток мозга, которые не восстанавливаются"),
        ("💭", "Память и концентрация", "Снижается способность запоминать, концентрироваться и принимать решения. У молодёжи эффект в 2 раза сильнее"),
        ("😔", "Психика", "Алкоголь — депрессант. Он усиливает тревогу и депрессию. Риск психических расстройств возрастает в 3–4 раза"),
        ("🔄", "Зависимость", "Мозг перестраивается: без алкоголя — «ломка», раздражительность, бессонница. Это физиологическая болезнь"),
        ("⚡", "Инсульт", "Даже умеренное употребление повышает риск инсульта в 1,5–3 раза. У людей 25–35 лет — главная причина инсульта"),
    ], W)
    add_note_box(slide, "Особая опасность для подростков: мозг формируется до 25 лет. Алкоголь наносит необратимый урон интеллекту и эмоциям.", W, 6.75)


def slide_health_body(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    W = prs.slide_width
    add_label(slide, "Слайд 5 — Влияние на органы", W)
    add_title(slide, "Алкоголь поражает весь организм", W)
    add_bullets(slide, [
        ("❤️", "Сердце", "Кардиомиопатия, аритмия, гипертония. Риск инфаркта выше в 2–4 раза. Риск внезапной остановки сердца"),
        ("🫁", "Печень", "Жировой гепатоз → гепатит → цирроз (необратимо). Риск рака печени выше в 5 раз"),
        ("🤢", "Поджелудочная", "Панкреатит, хроническая форма ведёт к сахарному диабету и онкологии"),
        ("🦷", "Желудок и кишечник", "Язвы, гастрит, нарушение всасывания веществ, рак пищевода и желудка"),
        ("🦴", "Кости и мышцы", "Остеопороз, атрофия мышечной ткани, постоянная слабость, высокий риск переломов"),
        ("🛡️", "Иммунитет", "Снижение защиты: инфекции протекают тяжелее, раны хуже заживают, онкологические риски растут"),
    ], W, font_size=13)


def slide_youth(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    W = prs.slide_width
    add_label(slide, "Слайд 6 — Алкоголь и молодёжь", W)
    add_title(slide, "Почему особенно опасно в юном возрасте", W)
    add_bullets(slide, [
        ("🧬", "Формирование мозга", "До 25 лет мозг активно развивается. Алкоголь нарушает этот процесс: снижается IQ, ухудшается память и контроль эмоций"),
        ("⚡", "Зависимость быстрее", "У подростков зависимость формируется в 2–4 раза быстрее. 1–2 года регулярного употребления могут привести к алкоголизму"),
        ("📚", "Учёба и карьера", "Снижение успеваемости, прогулы. Начавшие пить в школе в 2 раза реже заканчивают вуз"),
        ("👥", "Социальные последствия", "Конфликты в семье, потеря друзей, правонарушения. 70% подростковых правонарушений — под воздействием алкоголя"),
        ("💊", "Ворота к наркотикам", "Большинство наркозависимых начинали именно с алкоголя в подростковом возрасте"),
    ], W)
    add_note_box(slide, "Средний возраст первого употребления алкоголя в России — 13–14 лет. Один из самых низких показателей в Европе.", W, 6.75)


def slide_myths(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    W = prs.slide_width
    add_label(slide, "Слайд 7 — Мифы и правда", W)
    add_title(slide, "Популярные мифы об алкоголе", W)
    add_bullets(slide, [
        ("❌", "Миф: «Красное вино полезно для сердца»", ""),
        ("✅", "Факт", "Это миф алкогольной индустрии. Исследования 2019 г. (195 стран): любая доза увеличивает онкологический риск"),
        ("❌", "Миф: «Пиво — это не алкоголь»", ""),
        ("✅", "Факт", "Пиво содержит 4–8% этанола. По объёму — первое место в структуре потребления алкоголя в России"),
        ("❌", "Миф: «Алкоголь снимает стресс»", ""),
        ("✅", "Факт", "Алкоголь — депрессант ЦНС. Временно притупляет эмоции, но после трезвения стресс и тревога усиливаются"),
        ("❌", "Миф: «Алкоголь согревает»", ""),
        ("✅", "Факт", "Сосуды расширяются — тело быстрее теряет тепло. Алкоголь — причина гибели от переохлаждения"),
    ], W, font_size=13)


def slide_prevention(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    W = prs.slide_width
    add_label(slide, "Слайд 8 — Методы профилактики", W)
    add_title(slide, "Как защитить себя и близких", W)
    add_bullets(slide, [
        ("👤", "Личный уровень", "Осознать проблему, вести дневник употребления, ставить цели"),
        ("👤", "Замена привычки", "Спорт, творчество, волонтёрство — здоровые альтернативы"),
        ("👤", "Окружение", "Формировать круг общения без алкоголя, не бояться просить помощи"),
        ("👤", "Обратиться к врачу", "Зависимость — болезнь, которая лечится. Это сила, не слабость"),
        ("👨‍👩‍👧", "Ранний разговор", "Говорить с детьми о вреде алкоголя с 8–10 лет, до первых контактов со сверстниками"),
        ("👨‍👩‍👧", "Личный пример", "Родители — главный образец: дети копируют поведение взрослых"),
        ("📞", "Горячая линия", "8-800-200-0-200 — бесплатно, анонимно, круглосуточно"),
    ], W, font_size=13)


def slide_conclusion(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    W = prs.slide_width
    add_label(slide, "Слайд 9 — Заключение", W)
    add_title(slide, "Выводы", W)
    add_bullets(slide, [
        ("🚫", "Нет безопасной дозы", "Любое количество алкоголя наносит вред — доказано ВОЗ и исследованиями в 195 странах"),
        ("🧠", "Алкоголизм — болезнь", "Хроническое заболевание мозга, требующее лечения. Осуждение не помогает — нужна поддержка"),
        ("👶", "Особая опасность в юности", "Мозг подростка уязвим вдвойне: зависимость развивается быстрее, вред глубже и долгосрочнее"),
        ("📉", "Колоссальный ущерб", "Причина 200+ болезней, миллионов смертей, сломанных семей и огромного экономического урона"),
        ("🌱", "Выход есть", "Зависимость поддаётся лечению. Чем раньше обратиться — тем лучше прогноз"),
        ("💡", "Профилактика — лучшая защита", "Знание, здоровое окружение и занятость — надёжные инструменты против алкоголизма"),
    ], W)

    tb = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), W - Inches(1.2), Inches(0.7))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Выполнил: Илья Садомов, 9Д  |  Руководитель: Екатерина Геннадьевна"
    r.font.size = Pt(11)
    r.font.color.rgb = DIM


def handler(event: dict, context) -> dict:
    """Генерирует PPTX-файл презентации (9 слайдов) и возвращает его в base64."""
    if event.get('httpMethod') == 'OPTIONS':
        return {
            'statusCode': 200,
            'headers': {
                'Access-Control-Allow-Origin': '*',
                'Access-Control-Allow-Methods': 'GET, OPTIONS',
                'Access-Control-Allow-Headers': 'Content-Type',
            },
            'body': ''
        }

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_title(prs)
    slide_intro(prs)
    slide_stats_world(prs)
    slide_stats_russia(prs)
    slide_health_brain(prs)
    slide_health_body(prs)
    slide_youth(prs)
    slide_myths(prs)
    slide_prevention(prs)
    slide_conclusion(prs)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    encoded = base64.b64encode(buf.read()).decode('utf-8')

    return {
        'statusCode': 200,
        'headers': {
            'Access-Control-Allow-Origin': '*',
            'Content-Type': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'Content-Disposition': 'attachment; filename="alkogol_i_zdorovye.pptx"',
        },
        'body': encoded,
        'isBase64Encoded': True,
    }
